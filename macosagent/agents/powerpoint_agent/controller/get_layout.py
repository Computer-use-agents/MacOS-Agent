import json
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_ppt_elements_robust(ppt_path: str) -> List[Dict]:
    """
    完全错误处理版PPT元素提取
    
    参数:
        ppt_path: PPT文件路径
        
    返回:
        包含每页元素详细信息的字典列表
    """
    try:
        prs = Presentation(ppt_path)
        slides_data = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            try:
                slide_info = build_slide_info(slide, slide_idx)
                slides_data.append(slide_info)
            except Exception as e:
                print(f"处理幻灯片 {slide_idx} 时出错: {str(e)}")
                continue
        
        return slides_data
    except Exception as e:
        print(f"读取PPT文件时发生严重错误: {str(e)}")
        return []

def build_slide_info(slide, slide_idx: int) -> Dict:
    """构建单张幻灯片信息"""
    slide_info = {
        "slide_number": slide_idx,
        "slide_id": getattr(slide, 'slide_id', None),
        "slide_layout": safe_get_slide_layout_name(slide),
        "background": extract_background_info(slide),
        "shapes": [],
        "notes": extract_notes(slide),
        "placeholders": []
    }
    
    # 处理所有形状
    for shape in getattr(slide, 'shapes', []):
        try:
            shape_info = extract_shape_info(shape)
            if shape_info:
                slide_info["shapes"].append(shape_info)
        except Exception as e:
            print(f"幻灯片 {slide_idx} 形状处理出错: {str(e)}")
            continue
    
    # 处理占位符
    if hasattr(slide, 'placeholders'):
        for placeholder in slide.placeholders:
            try:
                placeholder_info = extract_placeholder_info(placeholder)
                if placeholder_info:
                    slide_info["placeholders"].append(placeholder_info)
            except Exception as e:
                print(f"幻灯片 {slide_idx} 占位符处理出错: {str(e)}")
                continue
    
    return slide_info

def safe_get_slide_layout_name(slide) -> str:
    """安全获取幻灯片布局名称"""
    try:
        return getattr(getattr(slide, 'slide_layout', None), 'name', 'unknown')
    except:
        return 'unknown'

def extract_shape_info(shape) -> Optional[Dict]:
    """提取形状信息，完全错误处理"""
    if not shape:
        return None
    
    try:
        shape_type = get_shape_type(shape)
        shape_info = {
            "id": getattr(shape, 'shape_id', None),
            "name": getattr(shape, 'name', ''),
            "type": shape_type,
            "type_id": getattr(shape, 'shape_type', None),
            "position": get_shape_position(shape),
            "fill": safe_extract_fill_info(shape),
            "line": safe_extract_line_info(shape),
            "text": safe_extract_text_info(shape),
            "image": safe_extract_image_info(shape, shape_type),
            "table": safe_extract_table_info(shape, shape_type),
            "chart": safe_extract_chart_info(shape, shape_type),
            "group": safe_extract_group_info(shape, shape_type),
            "auto_shape": safe_extract_auto_shape_info(shape, shape_type),
            "freeform": safe_extract_freeform_info(shape, shape_type),
            "connector": safe_extract_connector_info(shape, shape_type)
        }
        
        # 移除空值
        return {k: v for k, v in shape_info.items() if v is not None}
    except Exception as e:
        print(f"提取形状信息时出错: {str(e)}")
        return None

def get_shape_type(shape) -> str:
    """安全获取形状类型"""
    try:
        shape_type = getattr(shape, 'shape_type', None)
        if shape_type is None:
            return 'UNKNOWN'
        
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "PICTURE"
        elif shape_type == MSO_SHAPE_TYPE.TABLE:
            return "TABLE"
        elif shape_type == MSO_SHAPE_TYPE.CHART:
            return "CHART"
        elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return "AUTO_SHAPE"
        elif shape_type == MSO_SHAPE_TYPE.FREEFORM:
            return "FREEFORM"
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            return "GROUP"
        elif shape_type == MSO_SHAPE_TYPE.LINE or shape_type == MSO_SHAPE_TYPE.CONNECTOR:
            return "CONNECTOR"
        else:
            return str(shape_type)
    except:
        return 'UNKNOWN'

def get_shape_position(shape) -> Dict:
    """安全获取形状位置信息"""
    position = {
        "left": None,
        "top": None,
        "width": None,
        "height": None,
        "rotation": None
    }
    
    for attr in position.keys():
        try:
            val = getattr(shape, attr, None)
            if val is not None:
                position[attr] = val.pt if hasattr(val, 'pt') else float(val) if val else None
        except:
            continue
    
    return position

def safe_extract_fill_info(shape) -> Optional[Dict]:
    """安全提取填充信息"""
    if not hasattr(shape, 'fill'):
        return None
    
    try:
        fill = shape.fill
        if not fill:
            return None
        
        fill_info = {"type": str(getattr(fill, 'type', 'UNKNOWN'))}
        
        # 纯色填充
        if getattr(fill, 'type', None) == 1:  # 纯色填充
            if hasattr(fill, 'fore_color'):
                try:
                    fore_color = fill.fore_color
                    if fore_color:
                        fill_info["color"] = safe_rgb_to_hex(getattr(fore_color, 'rgb', None))
                        fill_info["transparency"] = getattr(fore_color, 'alpha', None)
                except:
                    pass
        
        # 渐变填充
        elif getattr(fill, 'type', None) == 2:  # 渐变填充
            if hasattr(fill, 'gradient_stops'):
                try:
                    stops = []
                    for stop in fill.gradient_stops:
                        stop_info = {
                            "position": getattr(stop, 'position', None),
                            "color": safe_rgb_to_hex(getattr(getattr(stop, 'color', None), 'rgb', None)),
                            "transparency": getattr(getattr(stop, 'color', None), 'alpha', None)
                        }
                        if any(stop_info.values()):
                            stops.append(stop_info)
                    if stops:
                        fill_info["gradient_stops"] = stops
                except:
                    pass
        
        # 图片或纹理填充
        elif getattr(fill, 'type', None) == 3:  # 图片或纹理填充
            if hasattr(fill, 'picture'):
                try:
                    picture = fill.picture
                    if picture and hasattr(picture, 'image'):
                        fill_info["image"] = {
                            "file_name": getattr(picture.image, 'filename', None),
                            "size": getattr(picture.image, 'size', None),
                            "format": getattr(picture.image, 'ext', None)
                        }
                except:
                    pass
        
        return fill_info if len(fill_info) > 1 else None
    except Exception as e:
        print(f"提取填充信息时出错: {str(e)}")
        return None

def safe_extract_line_info(shape) -> Optional[Dict]:
    """最终版线条信息提取，彻底解决_NoneFill问题"""
    if not hasattr(shape, 'line'):
        return None
    
    try:
        line = shape.line
        if not line:
            return None
        
        line_info = {}
        
        # 1. 提取线条基本属性
        if hasattr(line, 'width') and line.width is not None:
            try:
                line_info["width"] = line.width.pt if hasattr(line.width, 'pt') else float(line.width)
            except (AttributeError, ValueError):
                pass
        
        if hasattr(line, 'dash_style') and line.dash_style is not None:
            try:
                line_info["style"] = str(line.dash_style)
            except:
                pass
        
        if hasattr(line, 'cap_type') and line.cap_type is not None:
            try:
                line_info["cap_type"] = str(line.cap_type)
            except:
                pass
        
        # 2. 完全安全的线条颜色提取
        if hasattr(line, 'fill'):
            try:
                fill = line.fill
                # 关键修复：检查fill是否有有效的类型和前景色
                if hasattr(fill, 'type'):
                    fill_type = fill.type
                    # 只有solid(1)、gradient(2)或patterned(3)填充才有前景色
                    if fill_type in (1, 2, 3) and hasattr(fill, 'fore_color'):
                        fore_color = fill.fore_color
                        if fore_color and hasattr(fore_color, 'rgb') and fore_color.rgb:
                            hex_color = safe_rgb_to_hex(fore_color.rgb)
                            if hex_color:
                                line_info["color"] = hex_color
            except Exception:
                # 这里不再打印错误，因为_NoneFill是正常情况
                pass
        
        return line_info if line_info else None
    except Exception as e:
        # 记录其他类型的错误
        print(f"提取线条基本信息时出错: {str(e)}")
        return None

def safe_extract_text_info(shape) -> Optional[Dict]:
    """安全提取文本信息"""
    if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
        return None
    
    try:
        text_frame = shape.text_frame
        if not text_frame:
            return None
        
        text_info = {
            "text": getattr(text_frame, 'text', ''),
            "paragraphs": []
        }
        
        if hasattr(text_frame, 'paragraphs'):
            for para in text_frame.paragraphs:
                try:
                    para_info = {
                        "text": getattr(para, 'text', ''),
                        "alignment": str(getattr(para, 'alignment', None)),
                        "font": safe_extract_font_info(para)
                    }
                    text_info["paragraphs"].append(para_info)
                except:
                    continue
        
        return text_info if text_info["text"] or text_info["paragraphs"] else None
    except Exception as e:
        print(f"提取文本信息时出错: {str(e)}")
        return None

def safe_extract_font_info(paragraph) -> Optional[Dict]:
    """安全提取字体信息"""
    if not hasattr(paragraph, 'font'):
        return None
    
    try:
        font = paragraph.font
        if not font:
            return None
        
        font_info = {
            "name": getattr(font, 'name', None),
            "size": getattr(getattr(font, 'size', None), 'pt', None) if hasattr(getattr(font, 'size', None), 'pt') else None,
            "bold": getattr(font, 'bold', None),
            "italic": getattr(font, 'italic', None),
            "color": safe_rgb_to_hex(getattr(getattr(font, 'color', None), 'rgb', None)),
            "underline": getattr(font, 'underline', None)
        }
        
        return {k: v for k, v in font_info.items() if v is not None}
    except:
        return None

def safe_extract_image_info(shape, shape_type: str) -> Optional[Dict]:
    """安全提取图片信息"""
    if shape_type != "PICTURE" or not hasattr(shape, 'image'):
        return None
    
    try:
        image = shape.image
        if not image:
            return None
        
        image_info = {
            "file_name": getattr(image, 'filename', None),
            "size": getattr(image, 'size', None),
            "format": getattr(image, 'ext', None),
            "dpi": getattr(image, 'dpi', None)
        }
        
        return {k: v for k, v in image_info.items() if v is not None}
    except Exception as e:
        print(f"提取图片信息时出错: {str(e)}")
        return None

def safe_extract_table_info(shape, shape_type: str) -> Optional[Dict]:
    """安全提取表格信息"""
    if shape_type != "TABLE" or not hasattr(shape, 'has_table') or not shape.has_table:
        return None
    
    try:
        table = shape.table
        if not table:
            return None
        
        table_info = {
            "rows": getattr(table, 'rows', None).count if hasattr(getattr(table, 'rows', None), 'count') else None,
            "columns": getattr(table, 'columns', None).count if hasattr(getattr(table, 'columns', None), 'count') else None,
            "cells": []
        }
        
        if hasattr(table, 'rows'):
            for row in table.rows:
                if hasattr(row, 'cells'):
                    for cell in row.cells:
                        try:
                            cell_info = {
                                "text": getattr(cell, 'text', ''),
                                "row_span": getattr(cell, 'span_height', None),
                                "col_span": getattr(cell, 'span_width', None),
                                "fill": safe_extract_fill_info(cell),
                                "font": safe_extract_font_info(cell.text_frame.paragraphs[0]) if hasattr(cell, 'text_frame') and hasattr(cell.text_frame, 'paragraphs') and cell.text_frame.paragraphs else None
                            }
                            table_info["cells"].append({k: v for k, v in cell_info.items() if v is not None})
                        except:
                            continue
        
        return table_info if table_info["rows"] and table_info["columns"] else None
    except Exception as e:
        print(f"提取表格信息时出错: {str(e)}")
        return None

def safe_extract_chart_info(shape, shape_type: str) -> Optional[Dict]:
    """安全提取图表信息"""
    if shape_type != "CHART" or not hasattr(shape, 'has_chart') or not shape.has_chart:
        return None
    
    try:
        chart = shape.chart
        if not chart:
            return None
        
        chart_info = {
            "chart_type": str(getattr(chart, 'chart_type', None)),
            "title": getattr(chart, 'chart_title', {}).text if hasattr(chart, 'chart_title') and hasattr(chart.chart_title, 'text') else None
        }
        
        if hasattr(chart, 'plots'):
            plot = chart.plots[0] if len(chart.plots) > 0 else None
            if plot and hasattr(plot, 'categories'):
                chart_info["categories"] = [
                    getattr(cat, 'label', '') 
                    for cat in plot.categories 
                    if hasattr(cat, 'label')
                ]
            
            chart_info["series_count"] = len(getattr(chart, 'series', []))
            
            series_data = []
            for s in getattr(chart, 'series', []):
                try:
                    series_data.append({
                        "name": getattr(s, 'name', None),
                        "values": getattr(s, 'values', None),
                        "format": str(getattr(s, 'format', None))
                    })
                except:
                    continue
            if series_data:
                chart_info["series_data"] = series_data
        
        return {k: v for k, v in chart_info.items() if v is not None}
    except Exception as e:
        print(f"提取图表信息时出错: {str(e)}")
        return None

def safe_extract_group_info(shape, shape_type: str) -> Optional[Dict]:
    """安全提取组合形状信息"""
    if shape_type != "GROUP" or not hasattr(shape, 'shapes'):
        return None
    
    try:
        shapes = []
        for sub_shape in shape.shapes:
            try:
                sub_shape_info = extract_shape_info(sub_shape)
                if sub_shape_info:
                    shapes.append(sub_shape_info)
            except:
                continue
        
        return {
            "shapes_count": len(shapes),
            "shapes": shapes
        } if shapes else None
    except Exception as e:
        print(f"提取组合形状信息时出错: {str(e)}")
        return None

def safe_extract_auto_shape_info(shape, shape_type: str) -> Optional[Dict]:
    """安全提取自动形状信息"""
    if shape_type != "AUTO_SHAPE":
        return None
    
    try:
        auto_shape_info = {
            "shape_name": str(getattr(shape, 'auto_shape_type', None))
        }
        
        if hasattr(shape, 'adjustments'):
            try:
                adjustments = []
                for adj in shape.adjustments:
                    try:
                        adjustments.append(float(adj))
                    except:
                        continue
                if adjustments:
                    auto_shape_info["adjustments"] = adjustments
            except:
                pass
        
        return auto_shape_info if auto_shape_info["shape_name"] or "adjustments" in auto_shape_info else None
    except Exception as e:
        print(f"提取自动形状信息时出错: {str(e)}")
        return None

def safe_extract_freeform_info(shape, shape_type: str) -> Optional[Dict]:
    """安全提取自由形状信息"""
    if shape_type != "FREEFORM" or not hasattr(shape, 'vertices'):
        return None
    
    try:
        vertices = []
        for vertex in shape.vertices:
            try:
                vertices.append((float(vertex[0]), float(vertex[1])))
            except:
                continue
        
        return {"vertices": vertices} if vertices else None
    except Exception as e:
        print(f"提取自由形状信息时出错: {str(e)}")
        return None

def safe_extract_connector_info(shape, shape_type: str) -> Optional[Dict]:
    """安全提取连接线信息"""
    if shape_type != "CONNECTOR":
        return None
    
    try:
        connector_info = {
            "begin_connected": getattr(shape, 'begin_connected', False),
            "end_connected": getattr(shape, 'end_connected', False)
        }
        
        return connector_info if any(connector_info.values()) else None
    except Exception as e:
        print(f"提取连接线信息时出错: {str(e)}")
        return None

def extract_background_info(slide) -> Optional[Dict]:
    """提取幻灯片背景信息"""
    if not hasattr(slide, 'background') or not slide.background:
        return None
    
    try:
        background = slide.background
        fill = getattr(background, 'fill', None)
        if not fill:
            return None
        
        return safe_extract_fill_info(background)
    except Exception as e:
        print(f"提取背景信息时出错: {str(e)}")
        return None

def extract_notes(slide) -> Optional[Dict]:
    """提取幻灯片备注信息"""
    if not hasattr(slide, 'has_notes_slide') or not slide.has_notes_slide:
        return None
    
    try:
        notes_slide = slide.notes_slide
        if not hasattr(notes_slide, 'notes_text_frame') or not notes_slide.notes_text_frame:
            return None
        
        text_frame = notes_slide.notes_text_frame
        notes_info = {
            "text": getattr(text_frame, 'text', ''),
            "paragraphs": []
        }
        
        if hasattr(text_frame, 'paragraphs'):
            for para in text_frame.paragraphs:
                try:
                    para_info = {
                        "text": getattr(para, 'text', ''),
                        "font": safe_extract_font_info(para)
                    }
                    notes_info["paragraphs"].append(para_info)
                except:
                    continue
        
        return notes_info if notes_info["text"] or notes_info["paragraphs"] else None
    except Exception as e:
        print(f"提取备注信息时出错: {str(e)}")
        return None

def extract_placeholder_info(placeholder) -> Optional[Dict]:
    """提取占位符信息"""
    shape_info = extract_shape_info(placeholder)
    if not shape_info:
        return None
    
    try:
        shape_info.update({
            "placeholder_type": str(getattr(getattr(placeholder, 'placeholder_format', None), 'type', None)),
            "placeholder_idx": getattr(getattr(placeholder, 'placeholder_format', None), 'idx', None)
        })
        return shape_info
    except Exception as e:
        print(f"提取占位符信息时出错: {str(e)}")
        return shape_info

def safe_rgb_to_hex(rgb_color) -> Optional[str]:
    """终极安全的RGB颜色转换"""
    if not rgb_color:
        return None
    try:
        # 处理各种可能的颜色对象
        if hasattr(rgb_color, '__class__') and rgb_color.__class__.__name__ == 'RGBColor':
            if all(hasattr(rgb_color, c) for c in ('r', 'g', 'b')):
                r = max(0, min(255, getattr(rgb_color, 'r', 0)))
                g = max(0, min(255, getattr(rgb_color, 'g', 0)))
                b = max(0, min(255, getattr(rgb_color, 'b', 0)))
                return "#{:02x}{:02x}{:02x}".format(r, g, b)
        return None
    except Exception:
        return None

# 使用示例
if __name__ == "__main__":
    ppt_path = "input.pptx"  # 替换为你的PPT文件路径
    slides_info = extract_ppt_elements_robust(ppt_path)
    
    # 打印第一页的信息（格式化输出）
    print(f"总页数: {len(slides_info)}")
    if slides_info:
        print("\n第一页信息:")
        print(json.dumps(slides_info[0], indent=2, ensure_ascii=False))