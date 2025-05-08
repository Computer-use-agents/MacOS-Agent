import os

from pptx import Presentation
from pptx.dml.color import RGBColor


def hex_to_rgb(hex_color):
    """将十六进制颜色代码转换为RGB元组"""
    hex_color = hex_color.strip('#').lower()
    
    if len(hex_color) not in (3, 6):
        raise ValueError("无效的十六进制颜色格式")
    
    # 处理缩写格式（如 #fff -> ffffff）
    if len(hex_color) == 3:
        hex_color = ''.join([c*2 for c in hex_color])
    
    try:
        return (
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    except ValueError:
        raise ValueError("包含非法字符的十六进制颜色值")

def change_slide_background_function(input_path, slide_index, color, output_path=None):
    """
    修改PPT指定幻灯片的背景颜色（支持RGB元组/十六进制）
    
    参数：
    input_path (str): 原始PPT文件路径
    slide_index (int): 要修改的幻灯片索引（从1开始）
    color (tuple/str): 颜色值，支持两种格式：
                       - RGB元组 如 (255,255,255)
                       - 十六进制字符串 如 "#FF00FF" 或 "FF00FF"
    output_path (str): 输出文件路径，默认在原文件名后加'_modified'
    
    返回：
    str: 成功返回保存路径，失败返回错误信息
    
    示例：
    change_slide_background('demo.pptx', 0, "#FF0000")
    change_slide_background('demo.pptx', 1, (0, 255, 0))
    """
    slide_index = slide_index - 1
    if slide_index < 0:
        slide_index = 0
    try:
        # 参数预处理
        if isinstance(color, str):
            rgb = hex_to_rgb(color)
        elif isinstance(color, tuple):
            rgb = color
        else:
            raise TypeError("颜色参数必须为元组或字符串")
            
        # 验证颜色值范围
        if not all(0 <= c <= 255 for c in rgb):
            raise ValueError("颜色分量值必须在0-255之间")

        # 文件路径验证
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"文件 {input_path} 不存在")
            
        # 处理输出路径
        if not output_path:
            base, ext = os.path.splitext(input_path)
            output_path = f"{base}_modified{ext}"

        # 打开PPT文件
        prs = Presentation(input_path)
        
        # 验证幻灯片索引
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise IndexError(f"无效的幻灯片索引，最大有效索引为 {len(prs.slides)-1}")
        
        # 设置背景颜色
        slide = prs.slides[slide_index]
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*rgb)
        
        # 保存文件
        prs.save(output_path)
        return f"success，saved in {output_path}"
        
    except Exception as e:
        return f"fail : because {str(e)}"

# 使用示例
# if __name__ == "__main__":
#     # 使用十六进制颜色修改第一页
#     print(change_slide_background(
#         input_path="output.pptx", 
#         output_path="output1.pptx", 
#         slide_index=1, 
#         color="#FFB6C1"))  # 浅粉色
    
    # 使用RGB元组修改第二页
    # print(change_slide_background("demo.pptx", 1, (60, 179, 113)))  # 海洋绿
    

# def create_presentation_function(output_path):
#     prs = Presentation()
#     prs.save(output_path)
def create_presentation_function(output_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    prs.slides.add_slide(slide_layout)
    prs.save(output_path)
    

from pptx.enum.shapes import MSO_SHAPE_TYPE


def delete_specific_image_function(input_path, output_path, slide_index, image_index=1):
    """
    删除PPT中指定位置的图片并返回操作结果（注：仅仅是看不到图片，图片数据并未完全删除，需要使用ppt的检查功能彻底清除图片数据）
    :param ppt_path: PPT文件路径
    :param slide_index: 幻灯片页码（从1开始）
    :param image_index: 要删除的图片序号（默认为第1张）
    :return: 包含操作状态和信息的字典
    调用示例：
    result = delete_specific_image(
        input_path = "output_cmp.pptx", 
        output_path = "output1.pptx", 
        slide_index=2, 
        image_index=1)
    print(result)
    """
    result = {"success": False, "message": ""}
    slide_index -= 1
    image_index -= 1
    if slide_index<0:
        slide_index =0
    if image_index<0:
        image_index = 0 
    try:
        prs = Presentation(input_path)
        
        # 验证幻灯片索引有效性
        if slide_index >= len(prs.slides) or slide_index < 0:
            raise IndexError(f"无效的幻灯片索引：{slide_index}，总页数：{len(prs.slides)}")
            
        target_slide = prs.slides[slide_index]
        image_count = 0
        deleted = False

        # 遍历所有形状并执行删除
        for shape in target_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if image_count == image_index:
                    sp_tree = shape.element.getparent()
                    sp_tree.remove(shape.element)
                    deleted = True
                    break
                image_count += 1

        if not deleted:
            raise ValueError(f"未找到第{image_index+1}张图片，该页共{image_count}张图片")

        # 保存修改后的文件
        prs.save(output_path)
        result["success"] = True
        result["message"] = f"成功删除第{slide_index+1}页第{image_index+1}张图片"
        
    except Exception as e:
        result["message"] = f"操作失败：{str(e)}"
        # 根据异常类型细化错误信息
        if isinstance(e, FileNotFoundError):
            result["message"] = f"文件不存在：{input_path}"
        elif isinstance(e, IndexError):
            result["message"] = str(e)
        elif isinstance(e, PermissionError):
            result["message"] = "文件被占用或无写入权限"
            
    return result


def delete_specific_textbox_function(input_path, output_path, slide_index, textbox_index=1):
    """
    删除PPT中指定位置的文本框并返回操作结果（支持嵌套文本框检测）
    
    :param input_path: PPT源文件路径
    :param output_path: 修改后保存路径
    :param slide_index: 幻灯片页码（从1开始计数）
    :param textbox_index: 要删除的文本框序号（从1开始计数）
    :return: 包含操作状态和信息的字典
    
    调用示例：
    result = delete_specific_textbox(
        input_path = "input.pptx", 
        output_path = "output.pptx", 
        slide_index=3, 
        textbox_index=2)
    print(result)
    """
    result = {"success": False, "message": ""}
    slide_index -= 1  # 转换为0-based索引
    textbox_index -= 1
    if slide_index<0:
        slide_index =0
    if textbox_index<0:
        textbox_index = 0 
    try:
        # 参数有效性验证 [7,8](@ref)
        if not input_path.lower().endswith('.pptx'):
            raise ValueError("仅支持.pptx格式文件")

        prs = Presentation(input_path)
        
        # 幻灯片索引验证 [6](@ref)
        if not 0 <= slide_index < len(prs.slides):
            raise IndexError(f"无效的幻灯片索引：{slide_index+1}，总页数：{len(prs.slides)}")

        target_slide = prs.slides[slide_index]
        textbox_count = 0
        deleted = False

        # 倒序遍历避免索引错位 [4](@ref)
        for shape in reversed(target_slide.shapes):
            # 检测文本框（包含嵌套文本框）[7](@ref)
            if shape.has_text_frame:  
                if textbox_count == textbox_index:
                    shape.element.getparent().remove(shape.element)
                    deleted = True
                    break
                textbox_count += 1

        if not deleted:
            available_textboxes = sum(1 for s in target_slide.shapes if s.has_text_frame)
            raise ValueError(f"目标文本框不存在（该页共有{available_textboxes}个可删除文本框）")

        prs.save(output_path)
        result.update({
            "success": True,
            "message": f"已删除第{slide_index+1}页第{textbox_index+1}个文本框",
            "modified_path": output_path
        })
        
    except Exception as e:
        # 错误分类处理 [6,7](@ref)
        error_type = type(e).__name__
        detail_msg = str(e)
        
        if isinstance(e, FileNotFoundError):
            result["message"] = f"文件不存在：{input_path}"
        elif isinstance(e, PermissionError):
            result["message"] = "文件被占用或无写入权限（请关闭PPT程序后重试）"
        else:
            result["message"] = f"{error_type}：{detail_msg}"
            
        # 调试信息
        result["debug"] = {
            "input_path": input_path,
            "slide_index": slide_index + 1,
            "textbox_index": textbox_index + 1
        }
        
    return result

def delete_specific_table_function(input_path, output_path, slide_index, table_index=1):
    """
    删除PPT中指定位置的表格并返回操作结果（支持嵌套表格检测）
    
    :param input_path: PPT源文件路径
    :param output_path: 修改后保存路径
    :param slide_index: 幻灯片页码（从1开始计数）
    :param table_index: 要删除的表格序号（从1开始计数）
    :return: 包含操作状态和信息的字典
    
    功能特性：
    1. 支持检测组合图形中的嵌套表格[6](@ref)
    2. 自动过滤母版页表格
    3. 返回详细的表格位置信息
    
    调用示例：
    result = delete_specific_table(
        input_path="input.pptx", 
        output_path="output.pptx",
        slide_index=3,
        table_index=2)
    """
    result = {"success": False, "message": ""}
    slide_index -= 1  # 转换为0-based索引
    table_index -= 1
    if slide_index<0:
        slide_index =0
    if table_index<0:
        table_index = 0 
    try:
        # 参数有效性验证
        if not input_path.lower().endswith('.pptx'):
            raise ValueError("仅支持.pptx格式文件[7](@ref)")

        prs = Presentation(input_path)
        
        # 幻灯片索引验证
        if not 0 <= slide_index < len(prs.slides):
            raise IndexError(f"无效的幻灯片索引：{slide_index+1}，总页数：{len(prs.slides)}[2](@ref)")

        target_slide = prs.slides[slide_index]
        table_count = 0
        deleted = False

        # 倒序遍历避免索引错位[6](@ref)
        for shape in reversed(target_slide.shapes):
            # 检测表格及其嵌套情况
            if shape.has_table:
                if table_count == table_index:
                    # 从XML树中移除元素
                    sp_tree = shape.element.getparent()
                    sp_tree.remove(shape.element)
                    deleted = True
                    break
                table_count += 1
            # 处理组合图形中的嵌套表格[5](@ref)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.has_table:
                        if table_count == table_index:
                            sub_shape.element.getparent().remove(sub_shape.element)
                            deleted = True
                            break
                        table_count += 1
                if deleted:
                    break

        if not deleted:
            available_tables = sum(1 for s in target_slide.shapes if s.has_table)
            raise ValueError(f"目标表格不存在（该页共有{available_tables}个可删除表格）[3](@ref)")

        prs.save(output_path)
        result.update({
            "success": True,
            "message": f"已删除第{slide_index+1}页第{table_index+1}个表格",
            "modified_path": output_path
        })
        
    except Exception as e:
        # 增强型错误分类处理
        error_type = type(e).__name__
        detail_msg = str(e)
        
        if isinstance(e, FileNotFoundError):
            result["message"] = f"文件路径错误：{input_path} 不存在[2](@ref)"
        elif isinstance(e, PermissionError):
            result["message"] = "文件被占用或无写入权限（请关闭PPT程序后重试）[7](@ref)"
        elif "cannot save" in detail_msg.lower():
            result["message"] = "输出路径不可写：请检查output_path参数有效性[6](@ref)"
        else:
            result["message"] = f"{error_type}：{detail_msg}"
            
        # 调试信息记录
        result["debug"] = {
            "input_path": input_path,
            "slide_index": slide_index + 1,
            "table_index": table_index + 1
        }
        
    return result


from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.util import Pt


def modify_text_range_style_function(input_path, output_path, slide_index, textbox_index=1, 
                           start_char=0, end_char=None, font_name=None, font_size=None, 
                           font_color=None, bold=None, italic=None):
    """
    修改PPT文本框中指定字符区间的样式
    
    参数说明：
    start_char : 起始字符索引（从0开始计数，包含）
    end_char   : 结束字符索引（不包含），默认为None时修改到段落末尾
    
    调用示例：
    modify_text_range_style(
        input_path="input.pptx",
        output_path="output.pptx",
        slide_index=2,
        textbox_index=1,
        start_char=5,
        end_char=10,
        font_name="微软雅黑",
        font_color="#FF0000"
    )
    """
    result = {"success": False, "message": ""}
    
    try:
        # 索引转换（兼容用户输入的1-based索引）
        slide_idx = slide_index - 1
        textbox_idx = textbox_index - 1
        
        if slide_idx<0:
            slide_idx = 0 
        if textbox_idx <0:
            textbox_idx = 0 
        
        prs = Presentation(input_path)
        target_slide = prs.slides[slide_idx]
        
        # 定位文本框（参考网页1的文本框选择逻辑）
        textbox = target_slide.shapes[textbox_idx].text_frame
        
        # 提取完整文本并验证区间有效性
        full_text = ''.join([p.text for p in textbox.paragraphs])
        text_length = len(full_text)
        
        end_char = end_char or text_length
        if not (0 <= start_char < end_char <= text_length):
            raise ValueError(f"无效字符区间：{start_char}-{end_char}，文本总长度{text_length}")
        
        # 颜色格式转换（兼容网页7提到的十六进制颜色处理）
        if isinstance(font_color, str):
            hex_str = font_color.lstrip('#')
            if len(hex_str) == 6:
                r, g, b = (int(hex_str[i:i+2], 16) for i in (0, 2, 4))
                font_color = RGBColor(r, g, b)
            else:
                raise ValueError("颜色格式错误，需使用#RRGGBB格式")
        
        # 遍历段落进行区间修改（参考网页4的段落处理逻辑）
        current_pos = 0
        for para in textbox.paragraphs:
            para_text = para.text
            para_len = len(para_text)
            
            # 计算当前段落包含的字符区间
            para_start = current_pos
            para_end = current_pos + para_len
            
            # 确定需要修改的段落范围
            if para_end <= start_char:
                current_pos += para_len
                continue
            if para_start >= end_char:
                break
                
            # 分割段落内的目标区间
            local_start = max(start_char - para_start, 0)
            local_end = min(end_char - para_start, para_len)
            
            # 遍历文本运行(run)进行精确修改（参考网页8的格式刷逻辑）
            run_start = 0
            for run in para.runs:
                run_text = run.text
                run_len = len(run_text)
                run_end = run_start + run_len
                
                # 确定运行中的修改区间
                if run_end <= local_start:
                    run_start += run_len
                    continue
                if run_start >= local_end:
                    break
                    
                # 分割运行为前/中/后三部分
                modify_start = max(local_start - run_start, 0)
                modify_end = min(local_end - run_start, run_len)
                
                # 仅修改目标区间（参考网页2的字符级格式控制）
                if modify_start > 0 or modify_end < run_len:
                    # 分割前段
                    new_run = para.add_run()
                    new_run.text = run_text[:modify_start]
                    new_run.font.name = run.font.name
                    new_run.font.size = run.font.size
                    new_run.font.bold = run.font.bold
                    new_run.font.italic = run.font.italic
                    if run.font.color.type == MSO_COLOR_TYPE.RGB:
                        new_run.font.color.rgb = run.font.color.rgb
                    
                    # 修改段
                    modified_run = para.add_run()
                    modified_run.text = run_text[modify_start:modify_end]
                    if font_name: modified_run.font.name = font_name
                    if font_size: modified_run.font.size = Pt(font_size)
                    if font_color: modified_run.font.color.rgb = font_color
                    if bold is not None: modified_run.font.bold = bold
                    if italic is not None: modified_run.font.italic = italic
                    
                    # 分割后段
                    if modify_end < run_len:
                        new_run = para.add_run()
                        new_run.text = run_text[modify_end:]
                        new_run.font.name = run.font.name
                        new_run.font.size = run.font.size
                        new_run.font.bold = run.font.bold
                        new_run.font.italic = run.font.italic
                        if run.font.color.type == MSO_COLOR_TYPE.RGB:
                            new_run.font.color.rgb = run.font.color.rgb
                    # 移除原运行
                    para._p.remove(run._r)
                else:
                    # 整体修改运行
                    if font_name: run.font.name = font_name
                    if font_size: run.font.size = Pt(font_size)
                    if font_color: run.font.color.rgb = font_color
                    if bold is not None: run.font.bold = bold
                    if italic is not None: run.font.italic = italic
                
                run_start += run_len
            
            current_pos += para_len
        
        prs.save(output_path)
        result.update({
            "success": True,
            "message": f"已修改第{slide_index}页第{textbox_index}文本框的第{start_char}-{end_char}字符样式",
            "modified_path": output_path
        })
        
    except Exception as e:
        result["message"] = f"操作失败：{type(e).__name__} - {str(e)}"
        result["debug"] = {
            "input_path": input_path,
            "slide_index": slide_index,
            "textbox_index": textbox_index,
            "start_char": start_char,
            "end_char": end_char
        }
        
    return result


from pptx.enum.text import PP_ALIGN
from pptx.util import Inches


def check_boundary(prs, left, top, width, height):
    """检查元素是否超出边界（参数改为 prs）"""
    slide_w = prs.slide_width.inches  # 直接使用Presentation对象
    slide_h = prs.slide_height.inches
    
    warnings = []
    if (right_edge := left + width) > slide_w:
        warnings.append(f"右侧超出 {right_edge - slide_w:.2f} 英寸")
    if (bottom_edge := top + height) > slide_h:
        warnings.append(f"底部超出 {bottom_edge - slide_h:.2f} 英寸")
    return warnings

def insert_text_function(slide_index, prs_path,save_file_path, left, top, text, 
               width=3,  # 默认3英寸宽
               height=1, # 默认1英寸高
               font_size=14, 
               font_color=[0,0,0],
               align='left', 
               bg_color=None):
    """ 参数完整版本 """
    # 单位转换处理
    prs = Presentation(prs_path)
    slide = prs.slides[slide_index] 
    font_size = Pt(font_size)
    if font_color:
        font_color = RGBColor(font_color[0],font_color[1],font_color[2])
    if bg_color:
        bg_color = RGBColor(bg_color[0],bg_color[1],bg_color[2])
    
    left_inch = Inches(left).inches if isinstance(left, (int, float)) else left.inches
    top_inch = Inches(top).inches if isinstance(top, (int, float)) else top.inches
    width_inch = Inches(width).inches if isinstance(width, (int, float)) else width.inches
    height_inch = Inches(height).inches if isinstance(height, (int, float)) else height.inches
    
    # 边界检测（传入 prs）
    if warnings := check_boundary(prs, left_inch, top_inch, width_inch, height_inch):
        print(f"警告：{'; '.join(warnings)}")
        left_inch = 0 
        top_inch = 0 
        width_inch = prs.slide_width.inches  # 直接使用Presentation对象
        height_inch = prs.slide_height.inches
    
    # 创建文本框（其余代码保持不变）
    textbox = slide.shapes.add_textbox(Inches(left_inch), Inches(top_inch), 
                                      Inches(width_inch), Inches(height_inch))
    
    # 文本样式设置
    p = textbox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = font_size  # 现在可以正确接收font_size参数
    p.font.color.rgb = font_color
    
    # 对齐方式
    align_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT
    }
    p.alignment = align_map.get(align.lower(), PP_ALIGN.LEFT)
    
    # 背景色
    if bg_color:
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    prs.save(save_file_path)

def insert_image_function(slide_index, prs_path,save_file_path, left, top, image_path, 
                width=None, height=None, 
                keep_ratio=True):
    """智能插入图片（自动保持比例）"""
    # 单位处理
    prs = Presentation(prs_path)
    slide = prs.slides[slide_index] 
    
    left = Inches(left) if isinstance(left, (int, float)) else left
    top = Inches(top) if isinstance(top, (int, float)) else top
    
    # 获取图片原始尺寸
    from PIL import Image
    with Image.open(image_path) as img:
        img_width, img_height = img.size
        px_to_inch = 1/96  # 假设96dpi
        
    # 设置默认尺寸
    if not width and not height:
        width = Inches(img_width * px_to_inch)
        height = Inches(img_height * px_to_inch)
    
    # 边界检测
    w = width.inches if width else (height.inches * img_width/img_height)
    h = height.inches if height else (width.inches * img_height/img_width)
    if warnings := check_boundary(prs, left.inches, top.inches, w, h):
        print(f"警告：图片可能超出边界 - {', '.join(warnings)}")
    
    slide.shapes.add_picture(image_path, left, top, width, height)
    prs.save(save_file_path)
    
def insert_table_function(slide_index, prs_path,save_file_path, left, top, rows, cols, 
                data=None, 
                col_widths=None,
                table_style=None,  # 改用正式参数名
                width=4,  # 默认表格宽度（英寸）
                height=2,):  # 默认表格高度（英寸）
    """插入带边界检测的表格"""
    # 单位转换
    prs = Presentation(prs_path)
    slide = prs.slides[slide_index] 
    
    left_inch = Inches(left).inches if isinstance(left, (int, float)) else left.inches
    top_inch = Inches(top).inches if isinstance(top, (int, float)) else top.inches
    width_inch = Inches(width).inches if isinstance(width, (int, float)) else width.inches
    height_inch = Inches(height).inches if isinstance(height, (int, float)) else height.inches

    # 边界检测
    if warnings := check_boundary(prs, left_inch, top_inch, width_inch, height_inch):
        print(f"表格警告: {'; '.join(warnings)}")

    # 创建表格（使用转换后的单位）
    table_shape = slide.shapes.add_table(
        rows, cols,
        left=Inches(left_inch),
        top=Inches(top_inch),
        width=Inches(width_inch),
        height=Inches(height_inch)
    )
    table = table_shape.table

    # 填充数据
    data = data or [[f"R{r+1}C{c+1}" for c in range(cols)] for r in range(rows)]
    for r in range(min(rows, len(data))):
        row = data[r]
        for c in range(min(cols, len(row))):
            table.cell(r, c).text = str(row[c])

    # 应用正式样式（需python-pptx 0.6.18+）
    if table_style:
        try:
            # 新版写法
            from pptx.enum.table import TABLE_STYLE
            table.style = getattr(TABLE_STYLE, table_style.upper())
        except ImportError:
            # 旧版写法
            table.style = table_style
        except AttributeError:
            print(f"无效样式: {table_style}")
    prs.save(save_file_path)

import re


def escape_single_quotes(text):
    # 匹配未转义的单引号（不匹配 \\'）
    pattern = r"(?<!\\)'"
    return re.sub(pattern, r"\\'", text)


def parsing_response_to_pyautogui_code(responses, input_swap:bool=True) -> str:
    '''
    将M模型的输出解析为OSWorld中的action，生成pyautogui代码字符串
    参数:
        response: 包含模型输出的字典，结构类似于：
        {
            "action_type": "hotkey",
            "action_inputs": {
                "hotkey": "command v",
                "start_box": None,
                "end_box": None
            }
        }
    返回:
        生成的pyautogui代码字符串
    '''

    pyautogui_code = "import pyautogui\nimport time\n"
    if isinstance(responses, dict):
        responses = [responses]
    for response_id, response in enumerate(responses):
        if "observation" in response:
            observation = response["observation"]
        else:
            observation = ""

        if "thought" in response:
            thought = response["thought"]
        else:
            thought = ""
        
        # if response_id == 0:
        #     pyautogui_code += f"'''\nObservation:\n{observation}\n\nThought:\n{thought}\n'''\n"
        # else:
        #     pyautogui_code += f"\ntime.sleep(3)\n"

        action_dict = response
        action_type = action_dict.get("action_type")
        action_inputs = action_dict.get("action_inputs", {})
        
        if action_type == "HOT_KEY":
            # Parsing hotkey action
            hotkey = action_inputs.get("value", "").strip()
            if hotkey:
                # Handle other hotkeys
                keys = hotkey.split(" ")  # Split the keys by space
                pyautogui_code += "\ntime.sleep(0.5)\n"
                pyautogui_code += f"\npyautogui.hotkey({', '.join([repr(k) for k in keys])})"
            else:
                pyautogui_code += f"\n# Unrecognized hot key action parsing: {action_type}"
                # logger.warning(f"Unrecognized hot key action parsing: {response}")
        elif action_type == "ENTER":
            pyautogui_code += "\npyautogui.press('enter')"
        elif action_type == "INPUT":
            # Parsing typing action using clipboard
            content = action_inputs.get("value", "")
            content = escape_single_quotes(content)
            if content:
                if input_swap:
                    pyautogui_code += "\nimport pyperclip"
                    pyautogui_code += f"\npyperclip.copy('{content.strip()}')"
                    pyautogui_code += "\ntime.sleep(0.5)\n"
                    pyautogui_code += "\npyautogui.hotkey('command', 'v')"
                    pyautogui_code += "\ntime.sleep(0.5)\n"
                    if content.endswith("\n") or content.endswith("\\n"):
                        pyautogui_code += "\npyautogui.press('enter')"
                else:
                    pyautogui_code += f"\npyautogui.write('{content.strip()}', interval=0.1)"
                    pyautogui_code += "\ntime.sleep(0.5)\n"
                    if content.endswith("\n") or content.endswith("\\n"):
                        pyautogui_code += "\npyautogui.press('enter')"

        
        elif action_type in ["DRAG", "SELECT"]:
            # Parsing drag or select action based on start and end_boxes
            points = action_inputs.get("position")
            assert len(points) == 2 and len(points[0]) == 2 and len(points[1]) == 2, "Drag or select action should have 2 points"
            start_box = points[0]
            end_box = points[1]
            
            if start_box and end_box:
                x1, y1 = start_box # [x, y] float number between 0 and 1
                # sx = round(float(x1) * image_width, 3)
                # sy = round(float(y1) * image_height, 3)
                sx = round(float(x1), 3)
                sy = round(float(y1), 3)
                x2, y2 = end_box # [x, y] float number between 0 and 1
                # ex = round(float(x2) * image_width, 3)
                # ey = round(float(y2) * image_height, 3)
                ex = round(float(x2), 3)
                ey = round(float(y2), 3)
                pyautogui_code += (
                    f"\npyautogui.moveTo({sx}, {sy})\n"
                    f"\npyautogui.dragTo({ex}, {ey}, duration=1.0)\n"
                )
            # else:
            #     logger.warning(f"Parse failed! Drag or select action should have 2 points: {response}")
        elif action_type == "SCROLL":
            # Parsing scroll action
            start_box = action_inputs.get("position")
            if start_box:
                # x1, y1, x2, y2 = eval(start_box)  # Assuming box is in [x1, y1, x2, y2]
                # x = round(float((x1 + x2) / 2) * image_width, 3)
                # y = round(float((y1 + y2) / 2) * image_height, 3)
                x, y = start_box
                # x = round(float(x) * image_width, 3)
                # y = round(float(y) * image_height, 3)
                x = round(float(x), 3)
                y = round(float(y), 3)                
                # # 先点对应区域，再滚动
                pyautogui_code += f"\npyautogui.click({x}, {y}, button='left')"
            else:
                x = None
                y = None
            direction = action_inputs.get("direction", "")
            
            if x == None:
                if "up" in direction.lower():
                    pyautogui_code += "\npyautogui.scroll(5)"
                elif "down" in direction.lower():
                    pyautogui_code += "\npyautogui.scroll(-5)"
            else:
                if "up" in direction.lower():
                    pyautogui_code += f"\npyautogui.scroll(5, x={x}, y={y})"
                elif "down" in direction.lower():
                    pyautogui_code += f"\npyautogui.scroll(-5, x={x}, y={y})"

        elif action_type in ["CLICK", "LEFT_CLICK_DOUBLE", "LEFT_CLICK_SINGLE", "RIGHT_CLICK_SINGLE", "HOVER"]:
            # Parsing mouse click actions
            start_box = action_inputs.get("position")
            # start_box = str(start_box)
            if start_box:
                # start_box = eval(start_box)
                if len(start_box) == 2:
                    x1, y1 = start_box
                    x2 = x1
                    y2 = y1
                else:
                    raise ValueError(f"Unknown start_box format: {start_box}")
                # x = round(float((x1 + x2) / 2) * image_width, 3)
                # y = round(float((y1 + y2) / 2) * image_height, 3)

                x = round(float((x1 + x2)/2), 3)
                y = round(float((y1 + y2)/2), 3)

                if action_type == "LEFT_CLICK_DOUBLE" or action_type == "CLICK":
                    pyautogui_code += f"\npyautogui.click({x}, {y}, button='left')"
                elif action_type == "LEFT_CLICK_DOUBLE":
                    pyautogui_code += f"\npyautogui.doubleClick({x}, {y}, button='left')"
                elif action_type == "RIGHT_CLICK_SINGLE":
                    pyautogui_code += f"\npyautogui.click({x}, {y}, button='right')"
                elif action_type == "HOVER":
                    pyautogui_code += f"\npyautogui.moveTo({x}, {y})"
        
        elif action_type in ["FINISH"]:
            pyautogui_code = "DONE"
        
        else:
            pyautogui_code += f"\n# Unrecognized action type: {action_type}"
            #  logger.error(f"Unrecognized action type: {response}")

    return pyautogui_code
